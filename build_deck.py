#!/usr/bin/env python3
"""Deck PropHero — 5 bloques. Diseño Hero Club (ref. Masterclass). 9 slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0x13,0x14,0x1B); PRIMARY=RGBColor(0x2F,0x56,0xF6); RED=RGBColor(0xE8,0x55,0x4E)
GREEN=RGBColor(0x1F,0xA9,0x6B); WHITE=RGBColor(0xFF,0xFF,0xFF); CARD=RGBColor(0xFF,0xFF,0xFF)
INK=RGBColor(0x1A,0x1A,0x1A); RED_LEAD=RGBColor(0xD3,0x3A,0x33); BLUE_LEAD=RGBColor(0x2B,0x54,0xF5)
MUTED=RGBColor(0x8A,0x8F,0x9C); GREY=RGBColor(0x55,0x59,0x63); SUBC=RGBColor(0xB8,0xBD,0xC9)
EYE=RGBColor(0x9D,0xB4,0xFF); LOGO_BG=RGBColor(0x0E,0x0E,0x14); LOGO_LN=RGBColor(0x33,0x36,0x42)
FONT="Arial"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
FOOTER="PROPHERO · THE HERO CLUB   |   Cómo elegir la microzona perfecta · Junio 2026"

def bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
def glow(s):
    for x,y,w,h,c in [(-3.0,6.55,8.0,3.6,RGBColor(0x33,0x24,0x5E)),
                      (3.5,6.9,7.0,3.0,RGBColor(0x36,0x22,0x4C)),
                      (8.0,6.55,9.0,3.6,RGBColor(0x1C,0x30,0x64))]:
        e=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(w),Inches(h))
        e.fill.solid(); e.fill.fore_color.rgb=c; e.line.fill.background(); e.shadow.inherit=False
def base(): s=prs.slides.add_slide(BLANK); bg(s); glow(s); return s
def rrect(s,l,t,w,h,fill,radius=0.10,linec=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    try: sp.adjustments[0]=radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if linec is None: sp.line.fill.background()
    else: sp.line.color.rgb=linec; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def srect(s,l,t,w,h,fill):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def oval(s,l,t,d,fill):
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,d,d)
    o.fill.solid(); o.fill.fore_color.rgb=fill; o.line.fill.background(); o.shadow.inherit=False; return o
def txt(s,l,t,w,h,anchor=None):
    tf=s.shapes.add_textbox(l,t,w,h).text_frame; tf.word_wrap=True
    if anchor: tf.vertical_anchor=anchor
    return tf
def run(p,text,size,color=INK,bold=False,italic=False):
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.color.rgb=color
    r.font.bold=bold; r.font.italic=italic; r.font.name=FONT; return r
def line(tf,text,size,color=INK,bold=False,italic=False,align=PP_ALIGN.LEFT,sa=4,first=False):
    p=tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(sa); run(p,text,size,color,bold,italic); return p
def accent(s,l,t,w,color=PRIMARY,h=Inches(0.07)): rrect(s,l,t,w,h,color,radius=0.5)
def logo(s,l,t):
    rrect(s,l,t,Inches(1.05),Inches(0.95),LOGO_BG,radius=0.18,linec=LOGO_LN)
    tf=txt(s,l,t+Inches(0.12),Inches(1.05),Inches(0.75)); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for w,c,b in [("The",MUTED,False),("Hero",WHITE,True),("Club",MUTED,False)]:
        line(tf,w,12,c,bold=b,align=PP_ALIGN.CENTER,sa=0,first=(w=="The"))
def footer(s): line(txt(s,Inches(0.5),Inches(7.06),Inches(11),Inches(0.35)),FOOTER,8,MUTED,first=True)
def topbar(s,number,title,color=PRIMARY):
    srect(s,Inches(0),Inches(0),Inches(13.333),Inches(1.02),color)
    line(txt(s,Inches(0.45),Inches(0),Inches(1.5),Inches(1.02),MSO_ANCHOR.MIDDLE),number,40,WHITE,bold=True,first=True)
    line(txt(s,Inches(2.0),Inches(0),Inches(11.0),Inches(1.02),MSO_ANCHOR.MIDDLE),title,26,WHITE,bold=True,first=True)
def bar_sub(s,text): line(txt(s,Inches(0.6),Inches(1.2),Inches(12.2),Inches(0.7)),text,15,SUBC,italic=True,first=True)
def banner(s,text,color,top=Inches(5.95)):
    rrect(s,Inches(0.6),top,Inches(12.13),Inches(0.95),color,radius=0.12)
    line(txt(s,Inches(1.0),top,Inches(11.3),Inches(0.95),MSO_ANCHOR.MIDDLE),text,16,WHITE,bold=True,first=True)
def footnote(s,text,top=Inches(5.82)):
    line(txt(s,Inches(0.5),top,Inches(12.3),Inches(0.5)),text,11.5,MUTED,italic=True,first=True)
def pill(s,l,t,w,text,color=PRIMARY,size=15):
    rrect(s,l,t,w,Inches(0.62),color,radius=0.5)
    line(txt(s,l,t,w,Inches(0.62),MSO_ANCHOR.MIDDLE),text,size,WHITE,bold=True,align=PP_ALIGN.CENTER,first=True)
    accent(s,Emu(int(l)+int(w)//2-int(Inches(1.0))),Emu(int(t)+int(Inches(0.47))),Inches(2.0),color=WHITE,h=Inches(0.025))
def data_table(s,rows,top=Inches(1.95),height=Inches(3.55),widths=(3.7,2.6,6.03),val_col=1,
               header_color=PRIMARY,val_color=PRIMARY):
    n=len(rows); cols=len(rows[0])
    tbl=s.shapes.add_table(n,cols,Inches(0.5),top,Inches(12.33),height).table
    for ci,w in enumerate(widths): tbl.columns[ci].width=Inches(w)
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            c=tbl.cell(ri,ci); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.14); c.margin_top=Inches(0.03); c.margin_bottom=Inches(0.03)
            c.fill.solid(); c.fill.fore_color.rgb=header_color if ri==0 else (CARD if ri%2 else RGBColor(0xEE,0xF1,0xFB))
            p=c.text_frame.paragraphs[0]; c.text_frame.word_wrap=True
            if ri==0: run(p,val,12.5,WHITE,bold=True)
            elif ci==val_col: run(p,val,13,val_color,bold=True)
            else: run(p,val,12,INK,bold=(ci==0))
def index_card(s,x,y,w,h,num,title,desc):
    rrect(s,x,y,w,h,CARD,radius=0.12)
    cd=Inches(0.72); cy=Emu(int(y)+(int(h)-int(cd))//2); cx=Emu(int(x)+int(Inches(0.26)))
    oval(s,cx,cy,cd,PRIMARY)
    line(txt(s,cx,cy,cd,cd,MSO_ANCHOR.MIDDLE),num,15,WHITE,bold=True,align=PP_ALIGN.CENTER,first=True)
    tx=Emu(int(x)+int(Inches(1.25)))
    tf=txt(s,tx,y,Emu(int(w)-int(Inches(1.5))),h,MSO_ANCHOR.MIDDLE)
    line(tf,title,15,BLUE_LEAD,bold=True,first=True,sa=4)
    line(tf,desc,11.5,GREY)
def card(s,x,y,w,h,bullets,lead_color=BLUE_LEAD):
    rrect(s,x,y,w,h,CARD,radius=0.08)
    tf=txt(s,Emu(int(x)+int(Inches(0.28))),y,Emu(int(w)-int(Inches(0.56))),h,MSO_ANCHOR.MIDDLE)
    for i,(a,b) in enumerate(bullets):
        p=tf.paragraphs[0] if i==0 and not tf.paragraphs[0].runs else tf.add_paragraph()
        p.space_after=Pt(9)
        run(p,a+(" " if b else ""),13,lead_color,bold=True)
        if b: run(p,b,13,INK)

# ============================ 1 · PORTADA ============================
s=base(); logo(s,Inches(0.6),Inches(0.55))
line(txt(s,Inches(0.6),Inches(2.45),Inches(11.5),Inches(0.4)),"PROPHERO · DS&AI",13,EYE,bold=True,first=True)
tf=txt(s,Inches(0.55),Inches(2.95),Inches(12),Inches(2.0))
line(tf,"Cómo Elegir la",40,WHITE,first=True,sa=0); line(tf,"Microzona Perfecta",56,WHITE,bold=True,sa=0)
line(txt(s,Inches(0.6),Inches(5.4),Inches(11.5),Inches(0.6)),
     "Metodología de análisis  ·  Ignacio de la Cuba (AI Engineer, Growth)",14,MUTED,first=True)

# ============================ 2 · ÍNDICE ============================
s=base()
line(txt(s,Inches(0.6),Inches(0.55),Inches(12),Inches(0.4)),"AGENDA",13,EYE,bold=True,first=True)
line(txt(s,Inches(0.6),Inches(0.95),Inches(12),Inches(0.8)),"¿Qué veremos hoy?",34,WHITE,bold=True,first=True)
accent(s,Inches(0.62),Inches(1.72),Inches(2.6))
items=[("01","Por qué empezar desde arriba","El sesgo de confirmación y el embudo macro → micro"),
       ("02","Los 5 datos que sí o sí debes buscar","Filtro PropHero: transacciones · venta · alquiler · yield · población"),
       ("03","Cómo leer el mercado inmobiliario","Variaciones de venta y renta · ratio de esfuerzo · yield"),
       ("04","Empleo e industria como proxy de demanda","Sectores ancla · infraestructura · empresas grandes · noticias"),
       ("05","Validación cualitativa: lo que los datos no dan","Visita in situ · inmobiliarias · regla de oro")]
y0=int(Inches(2.05)); ch=int(Inches(0.9)); step=ch+int(Inches(0.1))
for i,(num,t,d) in enumerate(items):
    index_card(s,Inches(0.6),Emu(y0+i*step),Inches(12.13),Emu(ch),num,t,d)
footer(s)

# ============================ 3 · 01 EMPEZAR DESDE ARRIBA (SESGO) ============================
s=base()
topbar(s,"01.","Empezar desde arriba: evita el sesgo")
bar_sub(s,"Si empiezas por la calle que te gusta, buscas datos que la justifiquen. Eso es sesgo de confirmación.")
pill(s,Inches(0.6),Inches(2.35),Inches(5.85),"De abajo a arriba · el error",color=RED)
pill(s,Inches(6.88),Inches(2.35),Inches(5.85),"De arriba a abajo · PropHero",color=PRIMARY)
card(s,Inches(0.6),Inches(3.2),Inches(5.85),Inches(2.2),
     [("Parte de una calle o barrio que gusta",""),
      ("→","busca datos que lo justifiquen"),
      ("Sesgo de confirmación:","empiezas con la respuesta")],lead_color=RED_LEAD)
card(s,Inches(6.88),Inches(3.2),Inches(5.85),Inches(2.2),
     [("8.131 municipios,","sin idea previa"),
      ("→","el dato filtra y descarta"),
      ("La microzona","se valida al final, en campo")],lead_color=BLUE_LEAD)
banner(s,"Empieza por los datos, no por la calle. El embudo va de 8.131 municipios a 1 microzona.",PRIMARY,top=Inches(5.65))
footer(s)

# ============================ 4 · 02 LOS 5 DATOS (FILTRO PROPHERO) ============================
s=base()
topbar(s,"02.","Los 5 datos que sí o sí debes buscar")
bar_sub(s,"Filtro PropHero: 5 umbrales knock-out a nivel municipio. Pasa los 5 o queda fuera; después se rankea.")
data_table(s,[("Métrica","Umbral (knock-out)","Riesgo que mitiga"),
      ("Sale Transactions (12m)","≥ 30","Mercado sin liquidez: imposible salir"),
      ("Time to Sell","< 40 sem.","Capital atrapado en la venta"),
      ("Time to Rent","< 30 sem.","Vacancia prolongada, sin ingresos"),
      ("Yield bruto","≥ 6,5 %","Rentabilidad que no compensa el riesgo"),
      ("Crec. población (3 años)","≥ 0 %","Demanda estructural en declive")])
footnote(s,"Umbrales mínimos (suelo). Los benchmarks reales del cluster los superan con holgura.",top=Inches(5.7))
footer(s)

# ============================ 5 · 03 CÓMO LEER EL MERCADO ============================
s=base()
topbar(s,"03.","Cómo leer el mercado inmobiliario")
bar_sub(s,"Variaciones de venta y renta a 12m y 3 años: lo importante es la tendencia, no la foto fija.")
data_table(s,[("Métrica","Guía / objetivo","Cómo se evalúa"),
      ("Ratio de esfuerzo","< 30 % · > 40 % tensión","Alquiler anual ÷ renta media del hogar"),
      ("Yield bruto","≥ 6,5 % · España ~11,5 %","Renta anual ÷ precio de compra"),
      ("Absorción · time to rent","< 30 sem · España ~25 sem","Semanas hasta alquilar"),
      ("Precio de venta (€/m²)","Variación 12m y 3 años","Tendencia · España ~897 (2ª mano)"),
      ("Precio de alquiler","Variación 12m y 3 años","Recorrido de las rentas en el tiempo"),
      ("Stock de alquiler","Caída = tensión","Variación interanual de la oferta")],
      top=Inches(1.9),height=Inches(3.95),widths=(3.6,3.4,5.33))
footnote(s,"Medias de España (InvestMap · Brainsre) como termómetro; el ratio de esfuerzo es una guía orientativa.",top=Inches(5.95))
footer(s)

# ============================ 6 · 04 EMPLEO E INDUSTRIA (PROXY DEMANDA) ============================
s=base()
topbar(s,"04.","Empleo e industria como proxy de demanda")
bar_sub(s,"El empleo anticipa la demanda. Noticias, infraestructura y empresas grandes la sostienen.")
data_table(s,[("Indicador","Señal positiva","Alerta"),
      ("Empleadores","> 2 relevantes, sectores distintos","> 40 % en un solo empleador"),
      ("Sectores ancla","logística · agro · automóvil","cíclicos sin diversificar (turismo)"),
      ("Infraestructura","tren · carretera · proyectos futuros","sin conexión ni inversión"),
      ("Noticias (últimos 18 meses)","empresas que entran, obra pública","empresas que salen, cierres")],
      top=Inches(1.95),height=Inches(3.5),widths=(3.6,4.5,4.23))
footnote(s,"Cercanía a una empresa grande / subyacente = demanda sostenida que el dato de mercado no captura.",top=Inches(5.65))
footer(s)

# ============================ 7 · 05 VALIDACIÓN CUALITATIVA (VISITAS) ============================
s=base()
topbar(s,"05.","Validación cualitativa: lo que los datos no dan",GREEN)
bar_sub(s,"El municipio ya está validado. Ahora se baja a la microzona, en campo.")
rrect(s,Inches(0.6),Inches(1.8),Inches(12.13),Inches(0.8),LOGO_BG,radius=0.12,linec=GREEN,lw=1.5)
line(txt(s,Inches(1.0),Inches(1.8),Inches(11.3),Inches(0.8),MSO_ANCHOR.MIDDLE),
     "\"The street that on paper looks amazing — but it smells bad.\"",17,WHITE,bold=True,italic=True,first=True)
data_table(s,[("Acción en campo","Referencia","Por qué"),
      ("Visita in situ","corrige lag de 3–12 meses","El dato va por detrás del mercado"),
      ("Hablar con ≥ 2 inmobiliarias","producto, precio de cierre, calles","Absorción real, no de portal"),
      ("Observar la calle","comercio, carteles, viandantes","Empleo y demanda reales"),
      ("Regla de oro","campo > dato","Si discrepan, prevalece la visita")],
      top=Inches(2.75),height=Inches(2.85),widths=(4.3,3.0,5.03),
      header_color=GREEN,val_color=GREEN)
footnote(s,"La microzona sale de cruzar el dato con lo que cuentan los partners en campo.",top=Inches(5.72))
footer(s)

# ============================ 8 · CIERRE ============================
s=base()
tf=txt(s,Inches(0.6),Inches(1.3),Inches(11.8),Inches(1.8))
line(tf,"La microzona perfecta",36,WHITE,bold=True,first=True,sa=2)
line(tf,"no existe sin los datos correctos.",36,EYE,bold=True)
stats=[("5","Datos fundamentales (Filtro PropHero)"),
       ("8.131","Municipios como punto de partida"),
       ("1","Microzona validada en campo")]
sw,gap,sx=Inches(3.9),Inches(0.2),Inches(0.6)
for i,(num,lbl) in enumerate(stats):
    x=Emu(int(sx)+i*(int(sw)+int(gap)))
    rrect(s,x,Inches(3.85),sw,Inches(1.95),CARD,radius=0.10)
    tf=txt(s,x,Inches(4.0),sw,Inches(1.7),MSO_ANCHOR.MIDDLE)
    line(tf,num,44,PRIMARY,bold=True,align=PP_ALIGN.CENTER,first=True,sa=4)
    line(tf,lbl,12.5,INK,align=PP_ALIGN.CENTER)
footer(s)

# ============================ 9 · GRACIAS ============================
s=base()
line(txt(s,Inches(0),Inches(2.55),Inches(13.333),Inches(0.7)),"¿Preguntas?",26,MUTED,align=PP_ALIGN.CENTER,first=True)
line(txt(s,Inches(0),Inches(3.15),Inches(13.333),Inches(1.2)),"¡Gracias!",56,WHITE,bold=True,align=PP_ALIGN.CENTER,first=True)
logo(s,Inches(11.55),Inches(6.05))

out="/Users/ignaciodelacuba/Dev/claudio/PropHero_Microzona_v4.pptx"
prs.save(out); print("Guardado:",out,"·",len(prs.slides._sldIdLst),"slides")
